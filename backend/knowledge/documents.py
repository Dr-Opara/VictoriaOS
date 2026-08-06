from __future__ import annotations

import io
from pathlib import Path

from backend.core.logger import logger

_TEXT_EXTENSIONS = {".txt", ".md", ".csv"}
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff"}


class UnsupportedDocumentError(RuntimeError):
    """Raised when a file type has no text-extraction strategy."""


def extract_text(filename: str, content: bytes) -> str:
    """Extract plain text from an uploaded document.

    Dispatches by file extension. Supports plain text/markdown, PDF, Word
    (.docx), PowerPoint (.pptx), Excel (.xlsx), and images via OCR (when
    Tesseract is installed - see :func:`_extract_from_image`).
    """
    extension = Path(filename).suffix.lower()

    if extension in _TEXT_EXTENSIONS:
        return content.decode("utf-8", errors="replace")
    if extension == ".pdf":
        return _extract_from_pdf(content)
    if extension == ".docx":
        return _extract_from_docx(content)
    if extension == ".pptx":
        return _extract_from_pptx(content)
    if extension == ".xlsx":
        return _extract_from_xlsx(content)
    if extension in _IMAGE_EXTENSIONS:
        return _extract_from_image(content)

    raise UnsupportedDocumentError(f"No text extraction strategy for {extension!r} files.")


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150) -> list[str]:
    """Split text into overlapping chunks suitable for embedding.

    A simple fixed-size character splitter with overlap - sufficient for a
    personal knowledge base at this scale. Larger deployments would want
    sentence/paragraph-aware chunking; tracked as a roadmap refinement, not
    a blocker for a working RAG pipeline.
    """
    normalized = " ".join(text.split())
    if not normalized:
        return []

    chunks = []
    start = 0
    while start < len(normalized):
        end = start + chunk_size
        chunks.append(normalized[start:end])
        start = end - overlap if end - overlap > start else end

    return chunks


def _extract_from_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_from_docx(content: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(content))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _extract_from_pptx(content: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.append(shape.text_frame.text)

    return "\n".join(lines)


def _extract_from_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell is not None]
            if values:
                lines.append(" | ".join(values))

    return "\n".join(lines)


def _extract_from_image(content: bytes) -> str:
    """OCR an image via Tesseract, if it's installed on this system.

    ``pytesseract`` (the Python binding) is always installed as part of
    VictoriaOS's requirements, but it shells out to the real Tesseract OCR
    binary, which is a separate OS-level install
    (``apt install tesseract-ocr`` / ``brew install tesseract``). Without
    it, this raises a clear, actionable error rather than returning empty
    or fabricated text.
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError as error:
        raise UnsupportedDocumentError(
            "OCR requires the 'pytesseract' and 'Pillow' packages."
        ) from error

    try:
        image = Image.open(io.BytesIO(content))
        return pytesseract.image_to_string(image)
    except pytesseract.TesseractNotFoundError as error:
        logger.warning("OCR requested but the Tesseract binary is not installed.")
        raise UnsupportedDocumentError(
            "OCR requires the Tesseract OCR binary to be installed on this system "
            "(e.g. 'sudo apt install tesseract-ocr' / 'brew install tesseract'). "
            "The Python bindings are installed but the engine itself is missing."
        ) from error
